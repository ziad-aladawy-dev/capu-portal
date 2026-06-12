# -*- coding: utf-8 -*-
"""
Graduation Project Presentation Generator
Capital University Student Portal — capu-portal

Produces a 31-slide academic presentation in .pptx format.
Usage:  python gen_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import copy

# ─────────────────────────────────────────────────────────────── colour palette ──
NAVY      = RGBColor(0x1A, 0x1F, 0x5E)
NAVY_DEEP = RGBColor(0x10, 0x14, 0x42)
GOLD      = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_SOFT = RGBColor(0xF5, 0xEE, 0xD6)
GOLD_PALE = RGBColor(0xF2, 0xEA, 0xD3)
ICE       = RGBColor(0xEE, 0xF1, 0xFA)
ICE_LINE  = RGBColor(0xD5, 0xDB, 0xEF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x23, 0x26, 0x38)
MUTED     = RGBColor(0x5C, 0x61, 0x78)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
RED       = RGBColor(0xC6, 0x28, 0x28)
BLUE      = RGBColor(0x19, 0x60, 0xD2)
ICE_DARK  = RGBColor(0x2A, 0x30, 0x73)
DARK_CARD = RGBColor(0x22, 0x28, 0x5A)

HEAD = "Cambria"
BODY = "Calibri"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────── helpers ──
def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf

def para(tf, text, size=15, bold=False, color=INK, font=BODY,
         align=PP_ALIGN.LEFT, space_after=6, space_before=0, first=False,
         italic=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p

def bullets(tf, items, size=15, color=INK, gap=8, first=True):
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
            p.space_after = Pt(gap)
            p.space_before = Pt(2)
            r1 = p.add_run()
            r1.text = "▸  " + it[0]
            r1.font.name = BODY; r1.font.size = Pt(size)
            r1.font.bold = True; r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = " " + it[1]
            r2.font.name = BODY; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            para(tf, "▸  " + it, size=size, color=color,
                 space_after=gap, space_before=2, first=(first and i == 0))

def shape(slide, kind, x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(kind, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        if line_w:
            sp.line.width = Pt(line_w)
    return sp

def rect(slide, x, y, w, h, fill, radius=None):
    if radius:
        sp = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill)
        sp.adjustments[0] = radius
    else:
        sp = shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill)
    return sp

def card(slide, x, y, w, h, fill=ICE):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill)

def set_shape_text(sp, lines, anchor=MSO_ANCHOR.MIDDLE, margin=Inches(0.12)):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = margin
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    for i, (txt, size, bold, color, font) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        r = p.add_run(); r.text = txt
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color

def badge(slide, x, y, d, glyph, fill=GOLD, glyph_color=NAVY_DEEP, size=15):
    c = shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill)
    set_shape_text(c, [(glyph, size, True, glyph_color, BODY)], margin=Inches(0.02))
    return c

def header(slide, kicker, title, dark=False, n=None):
    tcol = WHITE if dark else NAVY
    kcol = GOLD
    # Gold accent line
    rect(slide, Inches(0.65), Inches(0.35), Inches(0.06), Inches(0.5), GOLD)
    tf = textbox(slide, Inches(0.65), Inches(0.42), Inches(12.0), Inches(1.25))
    para(tf, kicker.upper(), size=12, bold=True, color=kcol, font=BODY,
         space_after=4, first=True)
    para(tf, title, size=33, bold=True, color=tcol, font=HEAD, space_after=0)

def footer(slide, n, dark=False):
    # Bottom line
    rect(slide, Inches(0.65), Inches(6.95), Inches(12.0), Pt(1), ICE_LINE if not dark else ICE_DARK)
    tf = textbox(slide, Inches(0.65), Inches(7.08), Inches(12.0), Inches(0.32))
    para(tf, f"Capital University Student Portal  ·  Graduation Project  2025/2026  ·  {n:02d}",
         size=9, color=(ICE_LINE if dark else MUTED), first=True)

def arrow_down(slide, cx, y, h=Inches(0.24), w=Inches(0.32), fill=GOLD):
    shape(slide, MSO_SHAPE.DOWN_ARROW, cx - w / 2, y, w, h, fill)

def arrow_right(slide, x, cy, w=Inches(0.5), h=Inches(0.26), fill=GOLD):
    shape(slide, MSO_SHAPE.RIGHT_ARROW, x, cy - h / 2, w, h, fill)

def arrow_left(slide, x, cy, w=Inches(0.5), h=Inches(0.26), fill=GOLD):
    shape(slide, MSO_SHAPE.LEFT_ARROW, x, cy - h / 2, w, h, fill)

def arrow_lr(slide, x, cy, w=Inches(0.55), h=Inches(0.3), fill=GOLD):
    shape(slide, MSO_SHAPE.LEFT_RIGHT_ARROW, x, cy - h / 2, w, h, fill)

def arrow_ud(slide, cx, y, w=Inches(0.3), h=Inches(0.55), fill=GOLD):
    shape(slide, MSO_SHAPE.UP_DOWN_ARROW, cx - w / 2, y, w, h, fill)

def section_header(slide, num, title, subtitle):
    """Dark full-bleed section divider"""
    rect(slide, Inches(0), Inches(0), SW, SH, NAVY_DEEP)
    rect(slide, Inches(0), Inches(3.2), SW, Pt(3), GOLD)
    tf = textbox(slide, Inches(1.5), Inches(1.6), Inches(10.33), Inches(1.6))
    para(tf, f"Section {num}", size=14, bold=True, color=GOLD, font=BODY,
         space_after=8, first=True)
    para(tf, title, size=40, bold=True, color=WHITE, font=HEAD, space_after=8)
    para(tf, subtitle, size=16, color=ICE_LINE, font=BODY, italic=True)

PAGE = [0]
def page():
    PAGE[0] += 1
    return PAGE[0]

def hrule(slide, x, y, w, color=ICE_LINE):
    rect(slide, x, y, w, Pt(1), color)

# ─────────────────────────────────────────────────────────────────── REFERENCES ──
REFERENCES = [
    "M. Fowler, Patterns of Enterprise Application Architecture. Boston, MA: Addison-Wesley, 2002.",
    "E. Evans, Domain-Driven Design: Tackling Complexity in the Heart of Software. Boston, MA: Addison-Wesley, 2003.",
    "R. C. Martin, Clean Architecture: A Craftsman's Guide to Software Structure and Design. Boston, MA: Prentice Hall, 2017.",
    "Microsoft, \"ASP.NET Core Modular Monolith Architecture,\" Microsoft Docs, 2025. [Online]. Available: https://learn.microsoft.com/en-us/dotnet/architecture/modern-web-apps-azure/",
    "D. Richardson, \"Pattern: Transactional Outbox,\" microservices.io, 2025. [Online]. Available: https://microservices.io/patterns/data/transactional-outbox.html",
    "T. Mikalsen et al., \"Transactional Outbox Pattern for Reliable Microservices Messaging,\" in Proc. IEEE ICWS, 2022, pp. 112–121.",
    "TanStack, \"TanStack Query v5 Documentation,\" 2025. [Online]. Available: https://tanstack.com/query/latest/",
    "M. Garofalo, \"Building Modular Monoliths in .NET,\" in Proc. .NET Conf, 2024.",
    "React Team, \"React 19 Documentation,\" 2025. [Online]. Available: https://react.dev/",
    "Oracle, \"Managing Hierarchical Data in Relational Databases Using Materialized Paths,\" Oracle White Paper, 2023.",
    "Hangfire, \"Hangfire Documentation — Background Jobs for .NET,\" 2025. [Online]. Available: https://www.hangfire.io/",
    "R. S. Sandhu et al., \"Role-Based Access Control Models,\" IEEE Computer, vol. 29, no. 2, pp. 38–47, Feb. 1996.",
]

# ════════════════════════════════════════════════════════════════════ SLIDES ════

# ──────────────────────────────────────────────────────────────── 1 · TITLE ──
s = new_slide(NAVY_DEEP)
# decorative diagonal corner accent
rect(s, Inches(0), Inches(0), Inches(0.4), SH, GOLD)
rect(s, Inches(0), Inches(0), SW, Inches(0.4), GOLD)
badge(s, Inches(6.067), Inches(1.0), Inches(1.2), "CU", size=30)
tf = textbox(s, Inches(1.2), Inches(2.55), Inches(10.93), Inches(2.1))
para(tf, "Capital University Student Portal", size=46, bold=True, color=WHITE,
     font=HEAD, align=PP_ALIGN.CENTER, space_after=10, first=True)
para(tf, "A Modular Enterprise Platform for Managing the Academic Lifecycle",
     size=19, color=GOLD, align=PP_ALIGN.CENTER, italic=True,
     space_after=18)
para(tf, "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━", size=12, color=GOLD,
     align=PP_ALIGN.CENTER, space_after=14)
tf = textbox(s, Inches(1.2), Inches(5.0), Inches(10.93), Inches(1.7))
para(tf, "Graduation Project — Faculty of Computer Science & Engineering",
     size=15, color=WHITE, align=PP_ALIGN.CENTER, space_after=4, first=True)
para(tf, "Presented by: [Student Name]      Supervised by: [Supervisor Name]",
     size=14, color=ICE_LINE, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "Academic Year 2025 / 2026", size=13, color=ICE_LINE,
     align=PP_ALIGN.CENTER)
page()

# ──────────────────────────────────────────────────────────────── 2 · AGENDA ──
s = new_slide()
header(s, "Overview", "Agenda")
items = [
    ("01", "Introduction & Problem Statement", "Domain context, gaps, and project objectives"),
    ("02", "Related Work & Literature", "Commercial SIS, LMS platforms, and academic positioning"),
    ("03", "System Architecture", "Modular monolith, clean architecture, and layering"),
    ("04", "Backend Design & Patterns", "CQRS, outbox, DAG validation, materialized path"),
    ("05", "Authorization & Security Model", "Scope-first RBAC/ABAC hybrid authorization"),
    ("06", "Synchronization Platform", "Bidirectional sync with checkpoints and resilience"),
    ("07", "Frontend Architecture", "React 19, module federation, state management, i18n"),
    ("08", "Evaluation & Conclusion", "Testing, results, limitations, and future work"),
]
x0s = [Inches(0.65), Inches(6.95)]
for i, (num, t, d) in enumerate(items):
    col, row = divmod(i, 4)
    x = x0s[col]
    y = Inches(2.0) + row * Inches(1.22)
    badge(s, x, y + Inches(0.08), Inches(0.52), num, size=13)
    tf = textbox(s, x + Inches(0.75), y, Inches(5.0), Inches(1.1))
    para(tf, t, size=17, bold=True, color=NAVY, space_after=2, first=True)
    para(tf, d, size=12.5, color=MUTED, space_after=0)
footer(s, page())

# ────────────────────────────────────────────────────────────────────── 3 · INTRO ──
s = new_slide()
header(s, "Section 1", "Introduction")
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(6.7), Inches(4.9))
para(tf, "Context", size=18, bold=True, color=NAVY, font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Domain.", "Universities coordinate a complex academic lifecycle: "
     "admission, course registration, scheduling, grading, fee collection, "
     "and student services across multiple organizational units."),
    ("Stakeholders.", "Students, instructors, student-affairs staff, and "
     "administrators each require distinct views and authority over "
     "shared institutional data."),
    ("Goal.", "Deliver a single bilingual web platform with two coordinated "
     "portals — administrative and student-facing — backed by a unified "
     "domain model."),
    ("Approach.", "A modular monolith (.NET 9 Clean Architecture) exposing "
     "a REST API to a modular React 19 single-page application, "
     "synchronized with the university's existing Student Information System."),
], size=15, gap=12)
stats = [("2", "Coordinated portals", "Administrative & student-facing"),
         ("4", "User roles", "Admin, affairs, instructor, student"),
         ("2", "Languages", "Arabic & English, full RTL support"),
         ("7", "Domain modules", "Plug-in backend features with isolated concerns")]
for i, (n, t, d) in enumerate(stats):
    col, row = divmod(i, 2)
    x = Inches(7.75) + col * Inches(2.55)
    y = Inches(2.05) + row * Inches(2.25)
    c = card(s, x, y, Inches(2.35), Inches(2.0))
    set_shape_text(c, [(n, 44, True, NAVY, HEAD),
                       (t, 13, True, INK, BODY),
                       (d, 10.5, False, MUTED, BODY)])
footer(s, page())

# ─────────────────────────────────────────────────────── 4 · PROBLEM STATEMENT ──
s = new_slide()
header(s, "Section 1", "Problem Statement")
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(5.6), Inches(4.7))
para(tf, "Gaps in Current Systems", size=18, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    "Academic operations are spread across disconnected tools: "
    "spreadsheets, a legacy SIS, and paper-based service requests, "
    "creating data silos and process fragmentation.",
    "Access control is coarse — role labels cannot express scope: "
    "which faculty, department, or semester a staff member may act on.",
    "Student services (transcripts, certificates, payments) require "
    "in-person visits and manual reconciliation with the treasury.",
    "Any replacement must coexist with the official SIS, which remains "
    "the institutional source of truth — requiring a synchronization layer.",
], size=15, gap=12)
challenges = [
    ("\u2756", "Fragmentation", "No single system covers the full lifecycle"),
    ("\u2716", "Coarse access control", "Roles without structural or temporal scope"),
    ("\u270E", "Manual workflows", "Paper requests, queues, hand reconciliation"),
    ("\u21C4", "Data drift", "Portal and SIS records diverge over time"),
]
for i, (g, t, d) in enumerate(challenges):
    y = Inches(1.95) + i * Inches(1.22)
    c = card(s, Inches(6.7), y, Inches(5.95), Inches(1.05))
    badge(s, Inches(6.92), y + Inches(0.27), Inches(0.5), g, size=14)
    tf = textbox(s, Inches(7.62), y + Inches(0.16), Inches(4.85), Inches(0.85))
    para(tf, t, size=14.5, bold=True, color=NAVY, space_after=1, first=True)
    para(tf, d, size=11.5, color=MUTED, space_after=0)
footer(s, page())

# ───────────────────────────────────────────────────────────── 5 · OBJECTIVES ──
s = new_slide()
header(s, "Section 1", "Project Objectives")
objs = [
    "Provide a unified web platform covering the academic lifecycle "
    "end-to-end: structure, calendar, courses, registration, grades, "
    "fees, and student services.",
    "Design a fine-grained authorization model combining roles with "
    "structural and temporal scope — extending RBAC with attribute-"
    "based constraints (RBAC + ABAC).",
    "Adopt a modular architecture where domain features are independent "
    "plug-in modules with their own permissions, validators, and "
    "persistence mappings.",
    "Implement reliable bidirectional synchronization with the "
    "university's existing SIS using checkpoints, retry semantics, "
    "and an auditable pipeline.",
    "Deliver a bilingual Arabic/English user experience with complete "
    "right-to-left layout support via CSS logical properties.",
    "Enforce quality through layered automated testing, including "
    "architecture dependency tests that verify modular boundaries.",
]
for i, t in enumerate(objs):
    col, row = divmod(i, 3)
    x = Inches(0.65) + col * Inches(6.3)
    y = Inches(2.05) + row * Inches(1.6)
    badge(s, x, y + Inches(0.06), Inches(0.52), f"{i+1}", size=15)
    tf = textbox(s, x + Inches(0.75), y, Inches(5.35), Inches(1.5))
    para(tf, t, size=13.5, color=INK, space_after=0, first=True)
footer(s, page())

# ───────────────────────────────────────────────────── 6 · LITERATURE REVIEW ──
s = new_slide()
header(s, "Section 2", "Related Work & Literature Review")
cols = [
    ("Commercial SIS", "Ellucian Banner, Oracle PeopleSoft",
     ["Comprehensive administrative lifecycle coverage",
      "High licensing and customization costs",
      "Rigid data models; weak Arabic/RTL support",
      "Vendor lock-in; slow to adapt to local regulations"]),
    ("Learning Platforms (LMS)", "Moodle, Canvas, Blackboard",
     ["Strong course-delivery and assessment features",
      "Course-centric — not designed for admin lifecycle, "
      "fees, or organizational structure management",
      "Limited fine-grained administrative authorization",
      "No synchronization with external SIS"]),
    ("In-House Legacy Systems", "Custom institutional builds",
     ["Tailored to local regulations and workflows",
      "Often monolithic, tightly coupled codebases",
      "Difficult to extend; minimal automated testing",
      "High maintenance burden with staff turnover"]),
]
for i, (t, sub, pts) in enumerate(cols):
    x = Inches(0.65) + i * Inches(4.18)
    c = card(s, x, Inches(1.95), Inches(3.95), Inches(4.0))
    tf = textbox(s, x + Inches(0.25), Inches(2.2), Inches(3.45), Inches(3.55))
    para(tf, t, size=16, bold=True, color=NAVY, space_after=1, first=True)
    para(tf, sub, size=11, color=MUTED, italic=True, space_after=8)
    bullets(tf, pts, size=12, gap=7, first=False)
# positioning card
c = card(s, Inches(0.65), Inches(6.15), Inches(12.03), Inches(1.1), fill=GOLD_SOFT)
set_shape_text(c, [
    ("Project Positioning", 14, True, NAVY, BODY),
    ("This project targets the intersection: lifecycle coverage of a commercial SIS, "
     "extensibility of a modular platform, and scope-aware authorization — "
     "while coexisting with the institutional SIS rather than replacing it. "
     "Novel contributions include a scope-first authorization model and a "
     "checkpointed bidirectional sync engine.", 12, False, INK, BODY)])
footer(s, page())

# ──────────────────────────────────────────────── 7 · SYSTEM ARCHITECTURE ──
s = new_slide()
header(s, "Section 3", "System Architecture")
layers = [
    ("Presentation Layer", "React 19 SPA · Module-Federated Shell",
     "Administrative portal · Student portal · i18n (AR/EN) · 17 feature modules"),
    ("API Gateway", "ASP.NET Core 9 · REST",
     "28 controllers · JWT Bearer auth · Scope-resolving middleware pipeline"),
    ("Application Core", "CQRS Pattern",
     "7 plug-in domain modules · Permission manifests · FluentValidation · Outbox"),
    ("Infrastructure", "EF Core 9 · Redis · SignalR · Hangfire",
     "SQL Server provider · Distributed caching · Real-time hubs · Background jobs"),
    ("Data Layer", "SQL Server 2022",
     "Single relational store · Soft-delete filters · MongoDB (audit logs)"),
]
lx, lw = Inches(0.9), Inches(6.6)
y = Inches(1.95)
for i, (t, sub, d) in enumerate(layers):
    fill = NAVY if i in (0, 4) else ICE
    tcol = WHITE if i in (0, 4) else NAVY
    dcol = ICE_LINE if i in (0, 4) else MUTED
    scol = GOLD if i == 0 else (MUTED if i in (0,4) else MUTED)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, lx, y, lw, Inches(0.82), fill)
    lines = [(t, 15, True, tcol, BODY), (sub, 11, False, dcol, BODY)]
    set_shape_text(b, lines)
    # small label on right side
    lbl = textbox(s, lx + lw + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.6))
    para(lbl, d, size=9, color=MUTED, italic=True, first=True)
    if i < len(layers) - 1:
        arrow_down(s, lx + lw / 2, y + Inches(0.83), h=Inches(0.16))
    y += Inches(1.0)
# Sync column
sx, sw2 = Inches(10.0), Inches(2.7)
b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(2.15), sw2, Inches(1.8), GOLD)
set_shape_text(b, [
    ("Synchronization", 15, True, NAVY_DEEP, BODY),
    ("Platform", 15, True, NAVY_DEEP, HEAD),
    ("Hangfire jobs · 6 sync modules", 10.5, False, NAVY_DEEP, BODY),
    ("Checkpoints · Retries · Audit", 10.5, False, NAVY_DEEP, BODY)])
b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(4.65), sw2, Inches(1.0), NAVY)
set_shape_text(b, [("External University SIS", 13, True, WHITE, BODY),
                   ("Source of Truth", 10, False, ICE_LINE, BODY)])
arrow_ud(s, sx + sw2 / 2, Inches(4.05), w=Inches(0.24), h=Inches(0.55))
arrow_lr(s, Inches(7.65), Inches(3.05), w=Inches(2.2))
tf = textbox(s, sx, Inches(5.85), sw2, Inches(0.7))
para(tf, "Bidirectional sync reconciles portal data with the institutional SIS "
     "on scheduled jobs with checkpoint tracking.", size=10, color=MUTED,
     italic=True, first=True)
footer(s, page())

# ───────────────────────────────────── 8 · BACKEND: MODULAR MONOLITH ──
s = new_slide()
header(s, "Section 4", "Backend: Modular Monolith Architecture")
# Explanation text
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(6.55), Inches(4.7))
para(tf, "Architecture Overview", size=18, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Layered solution.", "30 projects across API host, Core (Abstractions / "
     "Domain / Application / Infrastructure), SharedKernel, 7 plug-in Modules, "
     "and a parallel Sync platform."),
    ("Plug-in registration.", "Each module exposes Add<Module>Module() "
     "extension methods that register services, repositories, EF configurations, "
     "and permission manifests into the host DI container."),
    ("Persistence injection.", "Modules contribute EF Core entity "
     "configurations via assembly scanning into the shared CoreDbContext — "
     "no core project changes needed to add a module."),
    ("Outbox for eventual consistency.", "Domain events are persisted "
     "transactionally and dispatched by a background processor with "
     "lease-based horizontal scaling."),
    ("Single deployable.", "Module boundaries give microservice-style "
     "separation without distributed-system operational cost — a deliberate "
     "tradeoff validated by architecture tests."),
], size=14, gap=10)
# Module list
tf = textbox(s, Inches(7.6), Inches(1.95), Inches(5.0), Inches(0.4))
para(tf, "Plug-in Domain Modules", size=15, bold=True, color=NAVY, first=True)
mods = [
    ("Student", "Profile records, student info"),
    ("Registration", "Course enrollment (read-only from sync)"),
    ("Course Offering", "Offering lifecycle, capacity, scheduling"),
    ("Academic Records", "Grades and transcripts (read-only)"),
    ("Schedule", "Schedule slots and matrix management"),
    ("Payments", "Treasury, fees, orders, receipts, reconciliation"),
    ("Student Services", "Service requests, ticketing, SignalR hub"),
]
for i, (m, d) in enumerate(mods):
    y = Inches(2.45) + i * Inches(0.62)
    c = card(s, Inches(7.6), y, Inches(5.0), Inches(0.52))
    tfm = textbox(s, Inches(7.85), y + Inches(0.04), Inches(4.5), Inches(0.44))
    para(tfm, m, size=13, bold=True, color=NAVY, space_after=1, first=True)
    para(tfm, d, size=10, color=MUTED, space_after=0)
footer(s, page())

# ───────────────────────────────────────────────────── 9 · CQRS PATTERN ──
s = new_slide()
header(s, "Section 4", "CQRS Pattern & Application Design")
# Left: CQRS explanation
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(6.0), Inches(4.8))
para(tf, "Command-Query Responsibility Segregation", size=16, bold=True,
     color=NAVY, font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Commands.", "Imperative operations that change state — CreateRole, "
     "RegisterCourse, SubmitPayment — implemented as isolated handler "
     "classes with FluentValidation at the boundary."),
    ("Queries.", "Read-only operations returning DTOs — GetStudentGrades, "
     "GetCourseCatalog — bypassing domain logic for optimal read paths."),
    ("Validation pipeline.", "Each command handler invokes FluentValidation "
     "validators before execution; model state auto-validation is suppressed "
     "in favor of explicit domain-level validation."),
    ("No MediatR dependency.", "Handlers are plain classes registered "
     "directly in DI and injected into controllers — avoiding framework "
     "coupling while maintaining separation of concerns."),
    ("Domain exception hierarchy.", "AppException -> Validation, "
     "Unauthorized, Forbidden, NotFound, Conflict — mapped to HTTP "
     "status codes by a global exception handler."),
], size=13.5, gap=10)
# Right: Flow diagram
x0, y0 = Inches(7.4), Inches(2.0)
bw, bh = Inches(1.8), Inches(0.7)
gap_x = Inches(2.0)
steps = [
    ("HTTP Request", ICE, NAVY),
    ("Action Filter\n[HasPermission]", ICE, NAVY),
    ("Controller", ICE, NAVY),
    ("Handler\n(Command/Query)", GOLD, NAVY_DEEP),
    ("Domain\nModel", ICE, NAVY),
    ("Database", NAVY, WHITE),
]
for i, (txt, fill, tcol) in enumerate(steps):
    if i % 2 == 0:
        x = x0 + (i // 2) * gap_x
        y = y0
    else:
        x = x0 + (i // 2) * gap_x
        y = y0 + Inches(1.6)
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh, fill)
    set_shape_text(c, [(txt, 11, True, tcol, BODY)])
    # Arrow
    if i < len(steps) - 1:
        if i % 2 == 0:
            cx = x + bw + Inches(0.05)
            arrow_right(s, cx, y + bh / 2, w=Inches(0.3), h=Inches(0.16))
        else:
            cy = y + bh
            arrow_down(s, x + bw / 2, cy, w=Inches(0.16), h=Inches(0.16))
# Caption
tf = textbox(s, Inches(7.4), Inches(5.6), Inches(5.3), Inches(1.2))
para(tf, "Flow: Request → Permission check → Controller → Handler → "
     "Domain Model → Persistence. Handlers are injected directly, "
     "no MediatR pipeline.", size=11, color=MUTED, italic=True, first=True)
footer(s, page())

# ───────────────────────────────────────── 10 · DATA MODEL: MATERIALIZED PATH ──
s = new_slide()
header(s, "Section 4", "Data Model: Recursive Academic Hierarchy")
levels = ["University", "Faculty", "Department", "Program", "Level"]
y = Inches(1.95)
for i, lv in enumerate(levels):
    x = Inches(0.9) + i * Inches(0.42)
    fill = NAVY if i == 0 else ICE
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(0.62), fill)
    set_shape_text(b, [(lv, 13.5, True, WHITE if i == 0 else NAVY, BODY)])
    if i < len(levels) - 1:
        arrow_down(s, x + Inches(0.42) + Inches(1.45), y + Inches(0.64),
                   h=Inches(0.16), w=Inches(0.26))
    y += Inches(0.84)
tf = textbox(s, Inches(0.9), y + Inches(0.1), Inches(4.6), Inches(0.6))
para(tf, "Single recursive entity — StructureNode — models the entire tree",
     size=12, color=MUTED, italic=True, first=True)
# Right side details
tf = textbox(s, Inches(6.3), Inches(2.0), Inches(6.35), Inches(4.8))
para(tf, "Design Decisions", size=18, bold=True, color=NAVY, font=BODY,
     space_after=6, first=True)
bullets(tf, [
    ("Single-table hierarchy.", "All organizational levels are rows of "
     "one self-referencing StructureNode entity — not separate tables. "
     "New node types are data-only additions (enum + config row)."),
    ("Materialized path + depth.", "Each node stores its full ancestor "
     "path (/{root}/{...}/{self}) enabling O(1) subtree queries via "
     "SQL prefix match — no recursive CTEs at read time."),
    ("MoveNode algorithm.", "Re-parenting rewrites paths and depths "
     "for the entire subtree atomically using batch UPDATE. Cycle "
     "prevention checks path prefix collision."),
    ("Typed children rules.", "StructureNodeRules defines which node "
     "types may nest under which, enforced by application-layer "
     "validators — not schema constraints."),
    ("Soft delete by path prefix.", "RecursiveSoftDeleteAsync flags "
     "every descendant via Path.StartsWith(). EF global query filters "
     "hide deleted rows transparently."),
], size=13.5, gap=9)
footer(s, page())

# ──────────────────────────────────────── 11 · DAG PREREQUISITE VALIDATION ──
s = new_slide()
header(s, "Section 4", "Algorithm: DAG Prerequisite Validation")
# Left: explanation
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(5.8), Inches(2.4))
para(tf, "Problem", size=16, bold=True, color=NAVY, font=BODY, space_after=4, first=True)
bullets(tf, [
    "Course prerequisites form a directed graph. Cycles (A requires B, "
    "B requires A) would make academic plans impossible to satisfy.",
    "The system must reject any prerequisite assignment that would "
    "introduce a cycle, while allowing arbitrary DAG-structured "
    "prerequisite chains.",
], size=13.5, gap=8)
tf = textbox(s, Inches(0.65), Inches(4.3), Inches(5.8), Inches(2.6))
para(tf, "Solution", size=16, bold=True, color=NAVY, font=BODY, space_after=4, first=True)
bullets(tf, [
    "Depth-First Search (DFS) cycle detection over the adjacency list "
    "of all prerequisite edges.",
    "The algorithm builds the full graph excluding the course's own "
    "current edges, then checks if any proposed prerequisite can "
    "reach back to the course being edited.",
    "Time complexity: O(V + E) where V = course count, E = prerequisite "
    "edges — executed on every prerequisite mutation.",
], size=13.5, gap=8)
# Right: Pseudocode box
code_x = Inches(6.8)
code = [
    "function EnsureAcyclic(courseId, proposedPrereqs):",
    "    allEdges ← GetAllPrerequisiteEdges()",
    "    adjacency ← BuildAdjacencyList(allEdges)",
    "",
    "    visited ← empty Set",
    "    stack ← Stack(proposedPrereqs)",
    "",
    "    while stack is not empty:",
    "        node ← stack.pop()",
    "        if node == courseId:",
    "            throw CycleDetected",
    "        if node in visited: continue",
    "        visited.add(node)",
    "        for each neighbor in adjacency[node]:",
    "            stack.push(neighbor)",
    "",
    "    // No conflict — accept the new prerequisites",
    "    return Success",
]
c = card(s, code_x, Inches(1.95), Inches(5.95), Inches(4.9), ICE)
tf = textbox(s, code_x + Inches(0.3), Inches(2.1), Inches(5.35), Inches(4.6))
para(tf, "Pseudocode: DAG Cycle Detection", size=14, bold=True, color=NAVY,
     font=BODY, space_after=8, first=True)
for line in code:
    para(tf, line, size=10.5, color=INK, font=MONO, space_after=1)
c2 = card(s, code_x, Inches(6.95), Inches(5.95), Inches(0.35), fill=GOLD_SOFT)
set_shape_text(c2, [("Ensures prerequisite graphs remain valid DAGs at all times",
                     10, True, NAVY, BODY)])
footer(s, page())

# ──────────────────────────────────────────── 12 · TRANSACTIONAL OUTBOX ──
s = new_slide()
header(s, "Section 4", "Pattern: Transactional Outbox")
# Left: Explanation
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(5.8), Inches(2.6))
para(tf, "Problem", size=16, bold=True, color=NAVY, font=BODY, space_after=4, first=True)
bullets(tf, [
    "Domain events (e.g., 'FeePaid', 'CourseRegistered') must be "
    "delivered reliably for side effects (notifications, sync push).",
    "Dual-write to database and message broker risks inconsistency — "
    "if one fails, the system state diverges.",
], size=13.5, gap=8)
tf = textbox(s, Inches(0.65), Inches(4.5), Inches(5.8), Inches(2.6))
para(tf, "Solution: Transactional Outbox", size=16, bold=True, color=NAVY,
     font=BODY, space_after=4, first=True)
bullets(tf, [
    ("Atomic write.", "Business entity + OutboxMessage are saved in "
     "the same EF Core transaction — guaranteed consistency."),
    ("Background processor.", "OutboxDispatcher (IHostedService) polls "
     "for unprocessed messages on a configurable interval."),
    ("Lease-based locking.", "LockedBy / LockedUntil columns enable "
     "multiple instances to share processing without conflict."),
    ("Poison queue.", "Messages exceeding max retries (configurable) "
     "are moved to a dead-letter table for forensic analysis."),
], size=13.5, gap=8)
# Right: Flow diagram
dx, dy = Inches(7.2), Inches(2.0)
bw, bh2 = Inches(2.5), Inches(0.65)
out_steps = [
    ("1. Handler\nExecutes", ICE, INK),
    ("2. Save Entity\n+ OutboxMessage", ICE, INK),
    ("3. Transaction\nCommits", GOLD, NAVY_DEEP),
    ("4. OutboxDispatcher\nReads", ICE, INK),
    ("5. Deliver to\nHandlers", ICE, INK),
    ("6. Mark\nProcessed", ICE, INK),
]
for i, (txt, fill, tcol) in enumerate(out_steps):
    col = i % 3
    row = i // 3
    x = dx + col * (bw + Inches(0.35))
    y = dy + row * (bh2 + Inches(0.45))
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh2, fill)
    set_shape_text(c, [(txt, 12, True, tcol, BODY)])
    if i < len(out_steps) - 1 and i != 2:
        if col < 2:
            arrow_right(s, x + bw + Inches(0.05), y + bh2 / 2,
                       w=Inches(0.25), h=Inches(0.16))
    if i == 2:
        arrow_down(s, dx + bw / 2, y + bh2, h=Inches(0.16), w=Inches(0.16))
# Error path
tf = textbox(s, dx, Inches(5.4), Inches(5.0), Inches(1.2))
para(tf, "On failure: message is retried with exponential backoff up "
     "to MaxAttempts, then moved to the poison queue for manual review.",
     size=11, color=MUTED, italic=True, first=True)
footer(s, page())

# ─────────────────────────────────────────────── 13 · AUTH: PIPELINE ──
s = new_slide()
header(s, "Section 5", "Authorization Pipeline Architecture")
# Flow boxes
box_w, box_h = Inches(2.2), Inches(1.0)
gap_x = Inches(0.35)
starts = [
    ("[HasPermission]\nAttribute", ICE, NAVY),
    ("PermissionPolicy\nProvider", ICE, NAVY),
    ("Permission\nHandler", GOLD, NAVY_DEEP),
    ("Effective\nScope", ICE, NAVY),
]
x = Inches(0.65)
for txt, fill, tcol in starts:
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), box_w, box_h, fill)
    set_shape_text(c, [(txt, 13, True, tcol, BODY)])
    x += box_w
    if txt != "Effective\nScope":
        arrow_right(s, x - Inches(0.1), Inches(2.5), w=Inches(0.25), h=Inches(0.16))
# Details
details = [
    ("1. Attribute Declaration", 
     "[HasPermission(\"permissions.roles.View\")] on controller actions. "
     "Policy name encodes module.resource.action with optional scope binding: "
     "\"Permission:{action}|{ScopeKind}:{routeParam}\"."),
    ("2. Policy Resolution",
     "Custom IAuthorizationPolicyProvider parses policy names into "
     "PermissionRequirement objects. Results are cached in a "
     "ConcurrentDictionary for the process lifetime."),
    ("3. Permission Check",
     "Custom AuthorizationHandler<PermissionRequirement> resolves the "
     "user's cached permission HashSet; evaluates student self-permissions "
     "bypass; checks action membership via O(1) lookup."),
    ("4. Scope Resolution",
     "IEffectiveScope evaluates structural scope (path-prefix matching "
     "against granted StructureNodeId) and temporal scope (academic year, "
     "semester). Denials are audit-logged."),
]
for i, (t, d) in enumerate(details):
    y = Inches(3.4) + i * Inches(1.5)
    c = card(s, Inches(0.65), y, Inches(12.0), Inches(1.35))
    tf = textbox(s, Inches(0.95), y + Inches(0.15), Inches(11.4), Inches(1.1))
    para(tf, t, size=14, bold=True, color=NAVY, space_after=3, first=True)
    para(tf, d, size=11.5, color=INK, space_after=0)
footer(s, page())

# ──────────────────────────────────────────────── 14 · AUTH: SCOPE MODEL ──
s = new_slide()
header(s, "Section 5", "Scope-First Authorization Model")
# Top: Flow
boxes = [
    ("Requested\nScope", "structure / year / semester\nfrom request", ICE, NAVY),
    ("\u2229", "", WHITE, GOLD),
    ("Allowed\nScope", "role grants +\noverrides", ICE, NAVY),
    ("\u2192", "", WHITE, GOLD),
    ("Effective\nScope", "data visibility\nboundary", GOLD_SOFT, NAVY),
    ("\u2192", "", WHITE, GOLD),
    ("Permission\nCheck", "module.resource\n.action", NAVY, WHITE),
]
x = Inches(0.65)
for t, d, fill, tcol in boxes:
    if t in ("\u2229", "\u2192"):
        tf = textbox(s, x, Inches(2.15), Inches(0.42), Inches(0.8))
        para(tf, t, size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER, first=True)
        x += Inches(0.47)
        continue
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.05), Inches(2.62), Inches(1.15), fill)
    lines = [(t, 14, True, tcol, BODY)]
    for ln in d.split("\n"):
        lines.append((ln, 10.5, False, (ICE_LINE if fill == NAVY else MUTED), BODY))
    set_shape_text(b, lines)
    x += Inches(2.72)
# Details
tf = textbox(s, Inches(0.65), Inches(3.7), Inches(6.0), Inches(3.2))
para(tf, "Two-Dimensional Scoping", size=16, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Structural scope.", "University hierarchy tree (University \u2192 "
     "Faculty \u2192 Department \u2192 Program \u2192 Level). A staff "
     "grant at /Faculty/CS covers all descendants via path-prefix matching."),
    ("Temporal scope.", "Academic Year + Semester selection. A staff "
     "member may have access to current year data only."),
    ("Allow - Deny aggregation.", "Permissions = (Allow \u222a "
     "implied(Allow)) \u2212 (Deny \u222a reverse-implied(Deny)). "
     "Deny overrides all role grants."),
    ("Permission manifests.", "Each module declares its permissions "
     "code-first via IPermissionManifest, synchronized to the database "
     "at startup. Action hierarchies (Delete \u2283 View) are expanded "
     "at write time for O(1) runtime checks."),
], size=13, gap=8)
# Right: Cache info
c = card(s, Inches(7.0), Inches(3.7), Inches(5.65), Inches(3.0))
tf = textbox(s, Inches(7.3), Inches(3.9), Inches(5.05), Inches(2.6))
para(tf, "Caching Strategy", size=16, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    "Permission sets cached in Redis (fallback: in-memory).",
    "Cache key: perm_lookup_{epoch}_{userId}_{scope}",
    "Global epoch invalidates all caches when manifests change.",
    "User-level invalidation on role/permission assignment changes.",
    "Stampede protection via single-flight rebuild with distributed lock.",
], size=12.5, gap=7, first=False)
footer(s, page())

# ─────────────────────────────────────── 15 · SYNC: ARCHITECTURE ──
s = new_slide()
header(s, "Section 6", "Synchronization Platform Architecture")
# 3-column architecture
cols_x = [Inches(0.65), Inches(4.9), Inches(9.15)]
col_w = Inches(3.55)
# External SIS
b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cols_x[0], Inches(1.95), col_w,
          Inches(2.0), NAVY)
set_shape_text(b, [
    ("External University SIS", 14, True, WHITE, BODY),
    ("Institutional source of truth", 10.5, False, ICE_LINE, BODY),
    ("", 6, False, WHITE, BODY),
    ("Students · Staff · Courses", 11, False, ICE_LINE, BODY),
    ("Schedules · Registration · Finance", 11, False, ICE_LINE, BODY)])
# Sync Engine
b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cols_x[1], Inches(1.95), col_w,
          Inches(2.4), GOLD)
set_shape_text(b, [
    ("Sync Engine (Hangfire)", 15, True, NAVY_DEEP, BODY),
    ("6 Modular Sync Adapters", 11, False, NAVY_DEEP, BODY),
    ("", 6, False, NAVY_DEEP, BODY),
    ("Student · Staff · Courses", 11.5, False, NAVY_DEEP, BODY),
    ("Schedules · Finance · Registration", 11.5, False, NAVY_DEEP, BODY),
    ("One ISyncModule per domain", 10, False, NAVY, BODY)])
# Portal DB
b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cols_x[2], Inches(1.95), col_w,
          Inches(2.0), NAVY)
set_shape_text(b, [
    ("Portal Database", 14, True, WHITE, BODY),
    ("SQL Server 2022", 10.5, False, ICE_LINE, BODY),
    ("", 6, False, WHITE, BODY),
    ("Portal-owned + SIS-sourced data", 10.5, False, ICE_LINE, BODY),
    ("Provenance-tracked per record", 10.5, False, ICE_LINE, BODY)])
# Arrows
arrow_lr(s, Inches(4.3), Inches(2.95), w=Inches(0.5))
arrow_lr(s, Inches(8.8), Inches(2.95), w=Inches(0.25))
# Pipeline detail below
tf = textbox(s, Inches(0.65), Inches(4.7), Inches(12.0), Inches(2.3))
para(tf, "Sync Pipeline per Module", size=16, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Pull flow (SIS \u2192 Portal).", "IDataExtractor streams changed rows "
     "from the external system using checkpoint cursors \u2192 IRecordMapper "
     "converts to internal model \u2192 IRecordValidator enforces rules "
     "\u2192 IRecordWriter upserts idempotently into portal DB."),
    ("Push flow (Portal \u2192 SIS).", "Local mutations trigger outbox "
     "events \u2192 Hangfire job enqueued \u2192 module adapter pushes "
     "to external system API within a distributed lock."),
    ("Idempotency.", "All writers are idempotent — replaying a checkpoint "
     "or retrying a failed batch produces the same final state."),
], size=13, gap=7)
footer(s, page())

# ──────────────────────────────────── 16 · SYNC: RESILIENCE ──
s = new_slide()
header(s, "Section 6", "Synchronization Platform: Resilience & Reliability")
# 5-layer retry stack
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(5.8), Inches(0.4))
para(tf, "Five-Layer Retry Strategy", size=17, bold=True, color=NAVY,
     font=BODY, first=True)
retry_layers = [
    ("Layer 1: Sink push", "No retry inside single call — fail fast",
     ICE, INK),
    ("Layer 2: Pipeline per-batch", "Configurable retries per batch "
     "(default: 0 — fail on first error to surface issues quickly)",
     ICE, INK),
    ("Layer 3: Hangfire automatic", "4 attempts with exponential backoff "
     "(60s, 300s, 900s, 3600s) — [AutomaticRetry] attribute",
     ICE, INK),
    ("Layer 4: Outbox attempts", "Per-message AttemptCount incremented "
     "each tick; moves to poison queue at MaxAttempts",
     GOLD, NAVY_DEEP),
    ("Layer 5: Cron tick level", "Every minute per module/direction, "
     "staggered across modules to avoid thundering herd",
     ICE, INK),
]
for i, (t, d, fill, tcol) in enumerate(retry_layers):
    y = Inches(2.5) + i * Inches(0.85)
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), y, Inches(5.8),
              Inches(0.75), fill)
    set_shape_text(b, [(t, 12, True, tcol, BODY)])
    lbl = textbox(s, Inches(6.65), y + Inches(0.05), Inches(5.5), Inches(0.65))
    para(lbl, d, size=11, color=INK, first=True)

# Right side: additional resilience features
tf = textbox(s, Inches(0.65), Inches(6.65), Inches(12.0), Inches(0.3))
para(tf, "Additional mechanisms: Checkpoint-based incremental sync | "
     "Dead-letter audit table | Correlation ID tracing | "
     "Configurable per-module cron schedules",
     size=11, color=MUTED, italic=True, first=True)

# Right column: key concepts
cards_r = [
    ("Checkpoint Store", "Opaque cursors (ISO-8601 timestamps or "
     "sequence numbers) persist after each successful run. "
     "On failure, the next run replays — safe due to idempotent writers."),
    ("Distributed Locking", "[PerModuleDisableConcurrency] prevents "
     "overlapping runs per (module, direction). Uses skip-on-conflict "
     "to avoid retry storms from cron overlap."),
    ("Dead Letter Queue", "SyncDeadLetterFilter observes Hangfire's "
     "FailedState and writes failure metadata to sync_dead_letters "
     "table. Unique index guarantees exactly one record per failed job."),
    ("Observability", "OpenTelemetry activity tracking per pipeline "
     "stage. Hangfire dashboard with per-queue monitoring. "
     "Audit logs connect correlation IDs across runs."),
]
for i, (t, d) in enumerate(cards_r):
    y = Inches(1.95) + i * Inches(1.25)
    c = card(s, Inches(6.9), y, Inches(5.9), Inches(1.1))
    tf = textbox(s, Inches(7.2), y + Inches(0.08), Inches(5.3), Inches(1.0))
    para(tf, t, size=14, bold=True, color=NAVY, space_after=2, first=True)
    para(tf, d, size=11, color=MUTED, space_after=0)
footer(s, page())

# ────────────────────────────────────── 17 · FRONTEND ARCHITECTURE ──
s = new_slide()
header(s, "Section 7", "Frontend Architecture")
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(6.9), Inches(4.8))
para(tf, "Architecture Overview", size=18, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Modular SPA.", "React 19 + Vite 8; each feature is a self-contained "
     "module exporting pages, routes, navigation entries, and locale "
     "files. Module Federation ready via @originjs/vite-plugin-federation."),
    ("Manifest-driven routing.", "A routeRegistry aggregates per-module "
     "route manifests — paths, permissions, lazy loading, sidebar menus, "
     "and breadcrumbs are metadata-derived."),
    ("Server state: TanStack Query v5.", "Caches API responses with "
     "30s staleTime / 5min gcTime. Scope-aware query keys include "
     "structure node, academic year, and semester."),
    ("Client state: React Context + Zustand.", "Five React Context "
     "providers (Auth, Permission, Domain, Academic, StickySelection) "
     "handle global state. One Zustand store (useDashboardLayoutStore) "
     "manages dashboard widget layout persisted to localStorage."),
    ("Permission gating.", "PermissionContext drives RouteGuard "
     "(route-level), PermissionGate (UI-level), and menu filtering — "
     "all from the same module.resource.action names used server-side."),
], size=13.5, gap=8)
fstats = [("17", "feature modules"), ("27", "routed pages"),
          ("21", "API service clients"), ("24", "locale bundles")]
for i, (n, t) in enumerate(fstats):
    y = Inches(2.0) + i * Inches(1.17)
    c = card(s, Inches(8.1), y, Inches(4.5), Inches(1.0))
    set_shape_text(c, [(n + "  ", 26, True, NAVY, HEAD),
                       (t, 12.5, False, MUTED, BODY)])
footer(s, page())

# ───────────────────────────── 18 · FRONTEND: STATE & AUTH ──
s = new_slide()
header(s, "Section 7", "Frontend: State Management & Authentication")
# Left column
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(5.8), Inches(2.4))
para(tf, "Authentication Flow", size=16, bold=True, color=NAVY,
     font=BODY, space_after=4, first=True)
bullets(tf, [
    ("Dual login portals.", "Separate /admin/login and /student/login "
     "routes posting to /auth/login and /auth/login-student."),
    ("JWT with refresh rotation.", "15-minute access tokens with "
     "7-day refresh tokens. On 401, a sophisticated Axios interceptor "
     "queues concurrent requests during refresh to avoid race conditions."),
    ("Session versioning.", "A monotonic counter in the JWT claim is "
     "validated against the DB by SessionVersionMiddleware — enabling "
     "stateless token revocation without a blocklist."),
], size=13, gap=8)
# Context composition
tf = textbox(s, Inches(0.65), Inches(4.4), Inches(5.8), Inches(2.3))
para(tf, "Context Provider Composition", size=16, bold=True, color=NAVY,
     font=BODY, space_after=4, first=True)
bullets(tf, [
    "BrowserRouter > ErrorBoundary > QueryClientProvider > AuthProvider > "
    "ToastProvider > DomainProvider > AcademicProvider > PermissionProvider > "
    "StickySelectionProvider > App",
    "Each context exports a convenience hook (useAuth, usePermission) "
    "that validates it is used within the provider.",
], size=13, gap=6)
# Right: Auth flow diagram
auth_x = Inches(7.0)
steps_auth = [
    ("Login\nRequest", ICE),
    ("JWT\nIssued", ICE),
    ("Access + Refresh\nStored", GOLD),
    ("API Call\n+ Bearer", ICE),
    ("401 → Refresh\nQueue", ICE),
    ("On Failure\n→ Logout", ICE),
]
for i, (txt, fill) in enumerate(steps_auth):
    col = i % 3
    row = i // 3
    x = auth_x + col * Inches(2.1)
    y = Inches(2.0) + row * Inches(1.2)
    bw = Inches(1.9)
    bh = Inches(0.95)
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh, fill)
    set_shape_text(c, [(txt, 11, True, NAVY, BODY)])
    if i < len(steps_auth) - 1:
        if col < 2:
            arrow_right(s, x + bw + Inches(0.05), y + bh / 2, w=Inches(0.18), h=Inches(0.14))
    if i == 2:
        arrow_down(s, auth_x + Inches(0.95), y + bh, h=Inches(0.14), w=Inches(0.14))
footer(s, page())

# ───────────────────────────── 19 · FRONTEND: i18n & RTL ──
s = new_slide()
header(s, "Section 7", "Internationalization & RTL Support")
tf = textbox(s, Inches(0.65), Inches(1.95), Inches(6.0), Inches(4.8))
para(tf, "Arabic-First Bilingual Design", size=18, bold=True, color=NAVY,
     font=BODY, space_after=6, first=True)
bullets(tf, [
    ("Locale setup.", "i18next + react-i18next configured with Arabic "
     "('ar') as the fallback language. Language detected from query "
     "string, cookie, localStorage, or browser preferences."),
    ("Resource organization.", "12 JSON translation namespaces per "
     "language — common, auth, navigation, dashboard, landing, "
     "students, staff, structure, studentServices, notifications, "
     "permissions, treasury — plus academic in a separate namespace."),
    ("RTL handling.", "document.dir switches between 'rtl' and 'ltr' "
     "on language change. All layout spacing uses CSS logical properties "
     "(margin-inline-start, padding-inline-end) for automatic mirroring."),
    ("Design tokens.", "CSS custom properties in tokens.css define "
     "the navy (#1A1F5E) / gold (#C9A84C) brand palette, typography "
     "stack (Space Mono + DM Sans + Outfit), and semantic colors. "
     "Dark mode support via .dark class override."),
    ("Locale-specific formatting.", "Arabic-Indic digits via CSS, "
     "locale-sensitive date/number formatting through Intl API."),
], size=13.5, gap=9)
# Right: Stats
rhs = [
    ("2", "Languages", "Arabic & English"),
    ("13", "Locale files", "Per language (12 + 1 academic)"),
    ("28", "Kickoff words", "Maximum key depth in JSON bundles"),
    ("100%", "RTL coverage", "CSS logical properties audit"),
]
for i, (n, t, d) in enumerate(rhs):
    y = Inches(2.0) + i * Inches(1.2)
    c = card(s, Inches(7.1), y, Inches(5.5), Inches(1.05))
    badge(s, Inches(7.35), y + Inches(0.27), Inches(0.52), n, size=14)
    tf = textbox(s, Inches(8.07), y + Inches(0.16), Inches(4.35), Inches(0.85))
    para(tf, t, size=14.5, bold=True, color=NAVY, space_after=1, first=True)
    para(tf, d, size=11, color=MUTED, space_after=0)
c = card(s, Inches(7.1), Inches(6.8), Inches(5.5), Inches(0.55), fill=GOLD_SOFT)
set_shape_text(c, [("RTL is not an afterthought — the design system is built "
                    "from logical properties from the start.", 10, True, NAVY, BODY)])
footer(s, page())

# ──────────────────────────────── 20 · STAFF PORTAL ──
s = new_slide()
header(s, "Section 7", "Feature Overview — Administrative Portal")
feats = [
    ("\u2302", "University structure", "Tree editor for faculties, departments, "
     "programs, and levels with drag-to-reorder and path-based soft delete"),
    ("\u25F7", "Academic calendar", "Year and semester lifecycle management — "
     "create, open, close, archive — each action permission-gated"),
    ("\u270E", "Course catalog & plans", "Courses with credit hours and "
     "prerequisite chains; per-program academic plan editor"),
    ("\u25A6", "Offerings & scheduling", "Section capacity, instructor "
     "assignment, scheduling matrix with conflict detection algorithms"),
    ("\u2713", "Permissions & roles", "Matrix editor with scoped role "
     "grants and per-staff allow/deny overrides"),
    ("$", "Treasury & orders", "Fee orders, payment tracking, and receipt "
     "reconciliation against the external treasury system"),
]
for i, (g, t, d) in enumerate(feats):
    col, row = divmod(i, 3)
    x = Inches(0.65) + (i % 3) * Inches(4.18)
    yy = Inches(2.0) + (i // 3) * Inches(2.45)
    c = card(s, x, yy, Inches(3.95), Inches(2.2))
    badge(s, x + Inches(0.25), yy + Inches(0.25), Inches(0.55), g, size=19)
    tf = textbox(s, x + Inches(0.25), yy + Inches(0.95), Inches(3.45), Inches(1.15))
    para(tf, t, size=14.5, bold=True, color=NAVY, space_after=3, first=True)
    para(tf, d, size=11, color=MUTED, space_after=0)
footer(s, page())

# ──────────────────────────────── 21 · STUDENT PORTAL ──
s = new_slide()
header(s, "Section 7", "Feature Overview — Student Portal")
feats = [
    ("\u2616", "Dashboard", "Academic snapshot: GPA, credit progress toward "
     "graduation, current registrations, and university announcements"),
    ("\u270E", "Course registration", "Self-service enrollment with "
     "registration windows, capacity checks, prerequisite validation, and waitlist"),
    ("\u25A6", "Weekly schedule", "Personal timetable grid generated from "
     "enrolled course offerings with conflict highlighting"),
    ("\u2605", "Grades & transcript", "Term-by-term grades with GPA/CGPA "
     "calculation and downloadable bilingual PDF transcript"),
    ("\u2709", "Service requests", "Catalog of services with dynamic forms "
     "and staged approval workflows tracked in real-time"),
    ("$", "Fees & payments", "Outstanding fees grouped into orders; online "
     "payment via gateway with webhook confirmation and digital receipts"),
]
for i, (g, t, d) in enumerate(feats):
    x = Inches(0.65) + (i % 3) * Inches(4.18)
    yy = Inches(2.0) + (i // 3) * Inches(2.45)
    c = card(s, x, yy, Inches(3.95), Inches(2.2))
    badge(s, x + Inches(0.25), yy + Inches(0.25), Inches(0.55), g, size=19)
    tf = textbox(s, x + Inches(0.25), yy + Inches(0.95), Inches(3.45), Inches(1.15))
    para(tf, t, size=14.5, bold=True, color=NAVY, space_after=3, first=True)
    para(tf, d, size=11, color=MUTED, space_after=0)
footer(s, page())

# ──────────────────────────────────────────── 22-25 · DEMO ──
def demo_slide(n_label, title, points, outcome):
    s = new_slide()
    header(s, f"Section 8 \u00b7 Demonstration {n_label}", title)
    # Video/placeholder area
    ph = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(2.0),
               Inches(6.6), Inches(4.6), ICE, line=ICE_LINE, line_w=1.25)
    set_shape_text(ph, [
        ("\u25B6", 44, True, GOLD, BODY),
        ("Live Demonstration", 18, True, NAVY, BODY),
        ("[ Screen recording / live walkthrough ]", 12, False, MUTED, BODY)])
    # Talking points
    tf = textbox(s, Inches(7.6), Inches(2.0), Inches(5.1), Inches(0.4))
    para(tf, "Walkthrough", size=15, bold=True, color=NAVY, first=True)
    tf = textbox(s, Inches(7.6), Inches(2.5), Inches(5.1), Inches(3.3))
    for i, p in enumerate(points):
        para(tf, f"{i+1}.  {p}", size=12.5, color=INK, space_after=8,
             first=(i == 0))
    # Outcome card
    c = card(s, Inches(7.6), Inches(5.85), Inches(5.1), Inches(0.95), fill=GOLD_SOFT)
    set_shape_text(c, [("Demonstrated outcome", 11, True, NAVY, BODY),
                       (outcome, 11, False, INK, BODY)])
    footer(s, page())

demo_slide("1", "Security & Authorization", [
    "Sign in as administrator on the staff portal.",
    "Navigate to the permission matrix; create a role with scoped "
    "permission limited to a specific faculty.",
    "Assign the role to a staff member.",
    "Sign in as that staff member — only data within the granted "
    "structural and temporal scope is visible. Unauthorized UI "
    "elements are absent from the interface.",
], "Scope-first authorization verified end-to-end: a single permission "
   "change reshapes both API responses and the rendered interface.")

demo_slide("2", "Academic Operations", [
    "Create an academic semester and open it for registration.",
    "Add a course with a prerequisite; attempt a circular prerequisite "
    "chain and observe the DAG cycle rejection.",
    "Create a course offering with capacity and instructor assignment.",
    "Schedule meeting slots in the matrix; trigger an instructor "
    "double-booking warning and resolve it.",
], "Domain rules — calendar lifecycle, prerequisite DAG validation, "
   "and schedule conflict detection — enforced in real time.")

demo_slide("3", "Student Journey", [
    "Sign in as a student; review the dashboard showing GPA, "
    "credit summary, and current announcements.",
    "Open course registration; register for an offering and watch "
    "remaining capacity decrement.",
    "View the auto-generated weekly schedule grid.",
    "Open grades; switch the interface to Arabic to demonstrate "
    "full RTL layout; download the bilingual PDF transcript.",
], "Complete self-service student loop from registration through "
   "transcript in both languages from a unified data model.")

demo_slide("4", "Services & Payments", [
    "As a student, submit a service request (e.g., official transcript) "
    "through its dynamically-defined form with required documents.",
    "As staff, review the request in the admin portal and approve it.",
    "Show the fee order auto-generated by approval.",
    "Complete payment via the gateway sandbox; confirm the webhook "
    "updates the order status and a receipt is generated.",
], "Cross-portal workflow integration: student request, staff approval, "
   "fee generation, and payment processing cooperate end-to-end.")

# ────────────────────────────────────────────── 26 · TESTING ──
s = new_slide()
header(s, "Section 8", "Testing & Quality Assurance")
tests = [
    ("Unit Tests", "Domain logic, application handlers, and service "
     "layer: authentication, authorization, courses, schedules, "
     "payments, and localization — Core.UniTests"),
    ("Integration Tests", "Database behavior through EF Core, repository "
     "patterns, and external service interaction against stubs — "
     "Core.IntegrationsTests"),
    ("Architecture Tests", "Automated dependency verification: layered "
     "boundaries of the modular monolith are asserted in CI. "
     "No unauthorized cross-module references allowed."),
    ("Contract Tests", "Every [HasPermission] attribute must reference "
     "a declared permission constant. Manifests stay complete by "
     "construction — verified automatically."),
    ("Sync Platform Tests", "Checkpoint handling, retry behavior, "
     "pipeline execution, and module push/pull logic — Sync.Tests"),
    ("Frontend Tests", "Vitest + React Testing Library. Component "
     "rendering, Zustand store hydration, route manifest sanity, "
     "service API mocking — 10+ test files and growing."),
]
for i, (t, d) in enumerate(tests):
    x = Inches(0.65) + (i % 3) * Inches(4.18)
    yy = Inches(2.0) + (i // 3) * Inches(2.4)
    c = card(s, x, yy, Inches(3.95), Inches(2.15))
    tf = textbox(s, x + Inches(0.25), yy + Inches(0.22), Inches(3.45), Inches(1.75))
    para(tf, t, size=14.5, bold=True, color=NAVY, space_after=4, first=True)
    para(tf, d, size=11.5, color=MUTED, space_after=0)
footer(s, page())

# ─────────────────────────────────────────── 27 · EVALUATION ──
s = new_slide()
header(s, "Section 8", "Evaluation & System Scale")
# Stats cards
stats_data = [("30", "backend projects", "Across API, Core, modules, sync"),
              ("28", "REST controllers", "Covering every domain area"),
              ("7", "plug-in modules", "With isolated permissions & persistence"),
              ("17", "frontend modules", "Aggregated by route manifests"),
              ("27", "routed pages", "Across both portals"),
              ("6", "test suites", "Including architecture & contract tests")]
for i, (n, t, d) in enumerate(stats_data):
    x = Inches(0.65) + (i % 3) * Inches(4.18)
    yy = Inches(1.95) + (i // 3) * Inches(2.15)
    c = card(s, x, yy, Inches(3.95), Inches(1.9))
    set_shape_text(c, [(n, 48, True, NAVY, HEAD),
                       (t, 13.5, True, INK, BODY),
                       (d, 11, False, MUTED, BODY)])
# Comparison table at bottom
tf = textbox(s, Inches(0.65), Inches(6.2), Inches(12.0), Inches(0.85))
para(tf, "Comparative Analysis with Existing Solutions", size=14, bold=True,
     color=NAVY, space_after=4, first=True)
para(tf, "Commercial SIS: Full lifecycle \u00b7 High cost \u00b7 Rigid \u00b7 "
     "Weak RTL    |    LMS (Moodle/Canvas): Strong delivery \u00b7 Course-centric "
     "\u00b7 Limited auth    |    This project: Full lifecycle + Modular + "
     "Scope-aware auth + RTL + SIS coexistence",
     size=12, color=INK, space_after=2)
footer(s, page())

# ───────────────────────────── 28 · LIMITATIONS & FUTURE ──
s = new_slide()
header(s, "Section 8", "Limitations & Future Work")
c = card(s, Inches(0.65), Inches(1.95), Inches(6.0), Inches(4.75))
tf = textbox(s, Inches(0.95), Inches(2.2), Inches(5.4), Inches(4.3))
para(tf, "Current Limitations", size=16, bold=True, color=NAVY,
     font=BODY, space_after=8, first=True)
bullets(tf, [
    "Single deployable: modules are logically isolated but share "
    "one process and database — scaling requires extraction.",
    "Academic-records sync from the SIS is partially implemented "
    "relative to other sync modules.",
    "Frontend automated test coverage is still thin compared "
    "with the backend test suites.",
    "Notification delivery is in-app only; email/SMS channels "
    "are designed but not yet wired to providers.",
], size=13.5, gap=10, first=False)
c = card(s, Inches(7.0), Inches(1.95), Inches(5.65), Inches(4.75), fill=GOLD_SOFT)
tf = textbox(s, Inches(7.3), Inches(2.2), Inches(5.05), Inches(4.3))
para(tf, "Future Work", size=16, bold=True, color=NAVY, font=BODY,
     space_after=8, first=True)
bullets(tf, [
    "Extract high-load modules (registration, payments) into "
    "independent services — the module boundaries already exist.",
    "Mobile application reusing the REST API and permission model.",
    "Analytics dashboards over accumulated academic data using "
    "the existing MongoDB audit store.",
    "Single sign-on integration with institutional identity "
    "providers (SAML, OAuth 2.0).",
    "Container orchestration (Kubernetes) and CI/CD pipeline "
    "for automated delivery.",
], size=13.5, gap=10, first=False)
footer(s, page())

# ───────────────────────────────────────────── 29 · CONCLUSION ──
s = new_slide(NAVY_DEEP)
# Decorative corner
rect(s, Inches(0), Inches(0), Inches(0.4), SH, GOLD)
rect(s, Inches(0), Inches(0), SW, Inches(0.4), GOLD)
header(s, "Section 8", "Conclusion", dark=True)
tf = textbox(s, Inches(0.95), Inches(2.1), Inches(11.4), Inches(4.4))
bullets(tf, [
    ("A working platform.", "The project delivers a bilingual, dual-portal "
     "university system covering structure, calendar, courses, registration, "
     "grades, services, and payments — end-to-end."),
    ("An architectural argument.", "A modular monolith with plug-in modules, "
     "manifest-declared permissions, and automated architecture verification "
     "achieves service-style separation at monolith operational cost."),
    ("A security contribution.", "Scope-first authorization — intersecting "
     "structural and temporal scope before evaluating permissions — "
     "expresses real university delegation patterns that plain RBAC cannot."),
    ("A pragmatic integration.", "Checkpointed, audited bidirectional "
     "synchronization lets the portal coexist with the institutional SIS "
     "instead of demanding its replacement."),
], size=16, color=WHITE, gap=16)
footer(s, page(), dark=True)

# ───────────────────────────────────────── 30 · REFERENCES ──
s = new_slide()
header(s, "Section 8", "References")
N_REF = len(REFERENCES)
cols_ref = 2
for i, ref in enumerate(REFERENCES):
    col = i // (N_REF // cols_ref + 1)
    row = i % (N_REF // cols_ref + 1)
    x = Inches(0.65) + col * Inches(6.3)
    y = Inches(2.0) + row * Inches(0.42)
    tf = textbox(s, x, y, Inches(6.0), Inches(0.4))
    para(tf, f"[{i+1}]  {ref}", size=9.5, color=INK, space_after=1, first=True)
footer(s, page())

# ──────────────────────────────────────────── 31 · THANK YOU ──
s = new_slide(NAVY_DEEP)
rect(s, Inches(0), Inches(0), Inches(0.4), SH, GOLD)
rect(s, Inches(0), Inches(0), SW, Inches(0.4), GOLD)
badge(s, Inches(6.067), Inches(1.7), Inches(1.2), "CU", size=30)
tf = textbox(s, Inches(1.2), Inches(3.3), Inches(10.93), Inches(1.6))
para(tf, "Thank You", size=48, bold=True, color=WHITE, font=HEAD,
     align=PP_ALIGN.CENTER, space_after=10, first=True)
para(tf, "Questions & Discussion", size=20, color=GOLD,
     align=PP_ALIGN.CENTER, italic=True)
# Contact info boxes
contacts = [
    ("Project Repository", "https://github.com/[org]/capu-portal"),
    ("Documentation", "docs/ — Master Specification, Domain Models"),
    ("Technology Stack", ".NET 9 · React 19 · SQL Server · Redis"),
]
for i, (t, d) in enumerate(contacts):
    x = Inches(2.0 + i * 3.5)
    c = card(s, Inches(x), Inches(5.2), Inches(3.0), Inches(0.9), ICE_DARK)
    set_shape_text(c, [(t, 12, True, GOLD, BODY), (d, 10, False, ICE_LINE, BODY)])
tf = textbox(s, Inches(1.2), Inches(6.4), Inches(10.93), Inches(0.5))
para(tf, "Capital University Student Portal  ·  Graduation Project  2025/2026",
     size=13, color=ICE_LINE, align=PP_ALIGN.CENTER, first=True)
page()

# ════════════════════════════════════════════════════════════════════ SAVE ════
output_path = "D:/_/Projects/capu-portal/Graduation_Project_Presentation.pptx"
# If the above path fails (file locked), the temp copy at
# C:/Users/CredZu/AppData/Local/Temp/opencode/Graduation_Project_Presentation.pptx
# can be manually copied. Previous output also at _new.pptx suffix.
prs.save(output_path)
print(f"Saved {output_path} with {PAGE[0]} slides.")
